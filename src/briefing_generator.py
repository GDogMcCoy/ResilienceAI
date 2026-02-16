"""
ResilienceAI - Executive Briefing Generator
Generates PDF and PPTX executive briefings for disaster vulnerability analysis.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import datetime
import pandas as pd
import numpy as np
from pathlib import Path
from config import PROCESSED_DIR, REPORTS_DIR

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                     TableStyle, PageBreak)
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class BriefingGenerator:
    """Generate executive briefings from county vulnerability data."""

    def __init__(self, df=None):
        if df is None:
            path = PROCESSED_DIR / "county_features.csv"
            if path.exists():
                self.df = pd.read_csv(path, dtype={"fips": str})
            else:
                self.df = None
        else:
            self.df = df

    def generate_county_brief(self, fips, output_format="pdf"):
        """Generate a single-county executive briefing."""
        if self.df is None:
            return {"error": "Data not loaded"}

        match = self.df[self.df["fips"] == str(fips)]
        if match.empty:
            return {"error": f"County {fips} not found"}

        county = match.iloc[0]

        if output_format == "pdf":
            return self._generate_pdf_brief(county)
        elif output_format == "pptx":
            return self._generate_pptx_brief(county)
        else:
            return self._generate_text_brief(county)

    def generate_state_brief(self, state_abbrev, output_format="pdf"):
        """Generate a state-level executive briefing."""
        if self.df is None:
            return {"error": "Data not loaded"}

        state_df = self.df[self.df["county_name"].str.contains(
            f", {state_abbrev}$", regex=True, na=False
        )]
        if state_df.empty:
            return {"error": f"No counties found for state {state_abbrev}"}

        if output_format == "pdf":
            return self._generate_state_pdf(state_abbrev, state_df)
        else:
            return self._generate_state_text(state_abbrev, state_df)

    def _generate_text_brief(self, county):
        """Generate a text-based executive briefing."""
        name = county.get("county_name", "Unknown County")
        risk = county.get("risk_score", 0)
        risk_level = county.get("risk_level", "Unknown")
        pop = county.get("total_population", 0)
        vuln = county.get("vulnerability_index", 0)
        iso = county.get("isolation_index", 0)
        disasters = county.get("disaster_count", 0)
        compound = county.get("compound_risk_count", 0)
        intervention = county.get("top_intervention", "N/A")
        redundancy = county.get("redundancy_score", 0)

        brief = f"""
{'='*60}
RESILIENCEAI EXECUTIVE BRIEFING
{name}
Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}
{'='*60}

RISK OVERVIEW
  Overall Risk Score: {risk:.3f} ({risk_level})
  Population: {pop:,}
  Vulnerability Index: {vuln:.3f}
  Isolation Index: {iso:.3f}
  Compound Risk Dimensions: {compound}/4

DISASTER HISTORY
  Total Disaster Declarations: {disasters}
  Acceleration: {'Increasing' if county.get('disaster_acceleration', 0) > 1 else 'Stable/Decreasing'}

INFRASTRUCTURE
  Redundancy Score: {redundancy:.3f}
  Zero Redundancy: {'YES - CRITICAL' if county.get('zero_redundancy_flag', 0) == 1 else 'No'}

RECOMMENDED INTERVENTION
  Top Priority: {intervention.replace('add_', 'Add ').replace('_', ' ').title()}
  Intervention Score: {county.get('top_intervention_score', 0):.3f}

KEY FINDINGS
"""
        findings = []
        if risk_level == "High":
            findings.append("- HIGH RISK: This county is in the top third of disaster vulnerability nationally.")
        if compound >= 3:
            findings.append(f"- COMPOUND RISK: County is high-risk across {compound} dimensions simultaneously.")
        if county.get("zero_redundancy_flag", 0) == 1:
            findings.append("- ZERO REDUNDANCY: Second-nearest hospital is >100km away. Single point of failure for healthcare.")
        if county.get("disaster_acceleration", 0) > 2.0:
            findings.append(f"- ACCELERATING DISASTERS: Disaster frequency has more than doubled in recent decade.")
        if county.get("poverty_pct", 0) > 20:
            findings.append(f"- HIGH POVERTY: {county.get('poverty_pct', 0):.1f}% poverty rate amplifies disaster impact.")
        if not findings:
            findings.append("- No critical alerts for this county.")

        brief += "\n".join(findings)
        brief += f"\n\n{'='*60}\nResilienceAI | MUIDSI 2026 | 100% Real Federal Data\n{'='*60}\n"

        # Save to file
        filename = f"briefing_{county.get('fips', 'unknown')}_{datetime.datetime.now().strftime('%Y%m%d')}.txt"
        output_path = REPORTS_DIR / filename
        with open(output_path, "w") as f:
            f.write(brief)

        return {"format": "text", "path": str(output_path), "content": brief}

    def _generate_pdf_brief(self, county):
        """Generate PDF executive briefing."""
        if not HAS_REPORTLAB:
            return self._generate_text_brief(county)

        name = county.get("county_name", "Unknown County")
        fips = county.get("fips", "unknown")
        filename = f"briefing_{fips}_{datetime.datetime.now().strftime('%Y%m%d')}.pdf"
        output_path = REPORTS_DIR / filename

        doc = SimpleDocTemplate(str(output_path), pagesize=letter,
                                topMargin=0.75*inch, bottomMargin=0.75*inch)
        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            "BriefTitle", parent=styles["Title"],
            textColor=HexColor("#4FC3F7"), fontSize=24,
            spaceAfter=6
        )
        subtitle_style = ParagraphStyle(
            "BriefSubtitle", parent=styles["Normal"],
            textColor=HexColor("#90A4AE"), fontSize=12,
            spaceAfter=20
        )
        heading_style = ParagraphStyle(
            "BriefHeading", parent=styles["Heading2"],
            textColor=HexColor("#4FC3F7"), fontSize=14,
            spaceBefore=12, spaceAfter=6
        )

        elements = []

        # Header
        elements.append(Paragraph("ResilienceAI Executive Briefing", title_style))
        elements.append(Paragraph(name, subtitle_style))
        elements.append(Paragraph(
            f"Generated: {datetime.datetime.now().strftime('%B %d, %Y')} | "
            f"FIPS: {fips}", styles["Normal"]
        ))
        elements.append(Spacer(1, 20))

        # Risk Overview Table
        elements.append(Paragraph("Risk Overview", heading_style))
        risk_data = [
            ["Metric", "Value"],
            ["Overall Risk Score", f"{county.get('risk_score', 0):.3f}"],
            ["Risk Level", str(county.get("risk_level", "Unknown"))],
            ["Population", f"{county.get('total_population', 0):,}"],
            ["Vulnerability Index", f"{county.get('vulnerability_index', 0):.3f}"],
            ["Isolation Index", f"{county.get('isolation_index', 0):.3f}"],
            ["Compound Risk Dimensions", f"{county.get('compound_risk_count', 0)}/4"],
            ["Disaster Declarations", str(county.get("disaster_count", 0))],
            ["Redundancy Score", f"{county.get('redundancy_score', 0):.3f}"],
        ]
        table = Table(risk_data, colWidths=[3*inch, 3*inch])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A1F2E")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#4FC3F7")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#546E7A")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#FFFFFF"), HexColor("#F5F5F5")]),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 15))

        # Recommended Intervention
        elements.append(Paragraph("Recommended Intervention", heading_style))
        intervention = county.get("top_intervention", "N/A")
        elements.append(Paragraph(
            f"<b>Priority Action:</b> {intervention.replace('add_', 'Add ').replace('_', ' ').title()}<br/>"
            f"<b>Intervention Impact Score:</b> {county.get('top_intervention_score', 0):.3f}",
            styles["Normal"]
        ))
        elements.append(Spacer(1, 15))

        # Footer
        elements.append(Paragraph(
            "ResilienceAI | MUIDSI 2026 | Built on 100% Real Federal Data",
            ParagraphStyle("Footer", parent=styles["Normal"],
                           textColor=HexColor("#90A4AE"), fontSize=8, alignment=1)
        ))

        doc.build(elements)
        return {"format": "pdf", "path": str(output_path)}

    def _generate_pptx_brief(self, county):
        """Generate PPTX executive briefing."""
        if not HAS_PPTX:
            return self._generate_text_brief(county)

        name = county.get("county_name", "Unknown County")
        fips = county.get("fips", "unknown")
        filename = f"briefing_{fips}_{datetime.datetime.now().strftime('%Y%m%d')}.pptx"
        output_path = REPORTS_DIR / filename

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "ResilienceAI"
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"Executive Briefing: {name}"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)
        p2.alignment = PP_ALIGN.CENTER

        # Risk Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Risk Overview"
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
        p.font.bold = True

        metrics = [
            ("Risk Score", f"{county.get('risk_score', 0):.3f}"),
            ("Risk Level", str(county.get('risk_level', 'N/A'))),
            ("Population", f"{county.get('total_population', 0):,}"),
            ("Vulnerability", f"{county.get('vulnerability_index', 0):.3f}"),
            ("Disasters", str(county.get('disaster_count', 0))),
            ("Redundancy", f"{county.get('redundancy_score', 0):.3f}"),
        ]

        for i, (label, value) in enumerate(metrics):
            col = i % 3
            row = i // 3
            x = Inches(0.5 + col * 4.2)
            y = Inches(1.5 + row * 2.5)

            txBox = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)

            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(36)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        prs.save(str(output_path))
        return {"format": "pptx", "path": str(output_path)}

    def _generate_state_pdf(self, state_abbrev, state_df):
        """Generate state-level PDF briefing."""
        if not HAS_REPORTLAB:
            return self._generate_state_text(state_abbrev, state_df)

        filename = f"briefing_state_{state_abbrev}_{datetime.datetime.now().strftime('%Y%m%d')}.pdf"
        output_path = REPORTS_DIR / filename

        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []

        title_style = ParagraphStyle(
            "StateTitle", parent=styles["Title"],
            textColor=HexColor("#4FC3F7"), fontSize=22
        )

        elements.append(Paragraph(
            f"ResilienceAI State Briefing: {state_abbrev}", title_style
        ))
        elements.append(Spacer(1, 20))

        # State summary stats
        elements.append(Paragraph("State Summary", styles["Heading2"]))
        summary_data = [
            ["Metric", "Value"],
            ["Counties", str(len(state_df))],
            ["Avg Risk Score", f"{state_df['risk_score'].mean():.3f}"],
            ["High Risk Counties", str((state_df['risk_level'] == 'High').sum())],
            ["Total Population", f"{state_df['total_population'].sum():,}"],
            ["Avg Poverty %", f"{state_df.get('poverty_pct', pd.Series([0])).mean():.1f}%"],
        ]
        table = Table(summary_data, colWidths=[3*inch, 3*inch])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A1F2E")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#4FC3F7")),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#546E7A")),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 15))

        # Top 10 highest risk counties
        elements.append(Paragraph("Top 10 Highest Risk Counties", styles["Heading2"]))
        top10 = state_df.nlargest(10, "risk_score")
        county_data = [["County", "Risk Score", "Population", "Top Intervention"]]
        for _, row in top10.iterrows():
            county_data.append([
                row.get("county_name", "?"),
                f"{row.get('risk_score', 0):.3f}",
                f"{row.get('total_population', 0):,}",
                str(row.get("top_intervention", "N/A")).replace("add_", "").replace("_", " ").title(),
            ])
        table = Table(county_data, colWidths=[2.5*inch, 1.5*inch, 1.5*inch, 1.5*inch])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A1F2E")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#4FC3F7")),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#546E7A")),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        elements.append(table)

        doc.build(elements)
        return {"format": "pdf", "path": str(output_path)}

    def _generate_state_text(self, state_abbrev, state_df):
        """Generate state-level text briefing."""
        brief = f"""
{'='*60}
RESILIENCEAI STATE BRIEFING: {state_abbrev}
Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}
{'='*60}

SUMMARY
  Counties: {len(state_df)}
  Avg Risk Score: {state_df['risk_score'].mean():.3f}
  High Risk Counties: {(state_df['risk_level'] == 'High').sum()}
  Total Population: {state_df['total_population'].sum():,}

TOP 10 HIGHEST RISK COUNTIES
"""
        top10 = state_df.nlargest(10, "risk_score")
        for _, row in top10.iterrows():
            brief += f"  {row.get('county_name', '?')}: {row.get('risk_score', 0):.3f} "
            brief += f"(pop: {row.get('total_population', 0):,})\n"

        brief += f"\n{'='*60}\n"

        filename = f"briefing_state_{state_abbrev}_{datetime.datetime.now().strftime('%Y%m%d')}.txt"
        output_path = REPORTS_DIR / filename
        with open(output_path, "w") as f:
            f.write(brief)
        return {"format": "text", "path": str(output_path), "content": brief}


if __name__ == "__main__":
    gen = BriefingGenerator()
    if gen.df is not None:
        fips = gen.df.iloc[0]["fips"]
        result = gen.generate_county_brief(fips, output_format="text")
        if "content" in result:
            print(result["content"][:500])
        print(f"\nSaved to: {result.get('path', 'N/A')}")
    else:
        print("Run pipeline first.")
